# -*- coding: utf-8 -*-
"""
住友商事ブランドPPTX ビルダー
template.pptx（表紙 / 目次 / 本文グリッド / エンドページ の4枚）を土台に、
JSON仕様（spec）からブランド準拠のPPTXを生成する。

使い方:
    python build.py <spec.json> <output.pptx> [--template <template.pptx>]

spec.json の形:
{
  "title": "プレゼンテーションタイトル",        # 表紙メインタイトル（必須）
  "subtitle": "サブタイトル／部署名",            # 表紙サブ（任意）
  "date": "住友商事株式会社 ｜ 2026年6月17日",   # 表紙下部の会社名・日付（任意）
  "toc": true,                                   # 目次スライドを入れるか（既定 true、本文2枚以上で自動）
  "thanks": "ご清聴いただきありがとうございました。",  # エンドの謝辞（任意）
  "slides": [                                    # 本文スライド（必須・1枚以上）
    {
      "section": "セクション名",                 # 上部の小見出し（ティール, 任意）
      "title": "ページタイトル",                 # 太字タイトル（必須）
      "lead": "リード文（1〜2行の説明）",         # 任意
      "body": [                                  # 本文。文字列 or {text,level,bold,bullet}
        "ただの段落",
        {"text": "箇条書き項目", "bullet": true},
        {"text": "字下げ項目", "bullet": true, "level": 1},
        {"text": "強調見出し", "bold": true}
      ],
      "figure": {                                # 図解（PowerPointネイティブ図形＝編集可）
        "viewBox": [1180, 500],                  # 図の座標系（全幅は約2.35:1）
        "shapes": [
          {"kind": "roundRect", "x": 40, "y": 60, "w": 300, "h": 120,
           "fill": "teal", "text": "入力", "size": 18, "bold": true},
          {"kind": "arrow", "points": [[350, 120], [420, 120]]},
          {"kind": "roundRect", "x": 430, "y": 60, "w": 300, "h": 120,
           "fill": "white", "line": "teal", "text": "処理", "size": 18},
          {"kind": "text", "x": 40, "y": 220, "w": 700, "h": 40,
           "text": "補足", "size": 16, "align": "l"}
        ]
      }
    }
  ]
}

設計ルール（ブランド分析に基づく）:
- 図解は必ず figure（オートシェイプ）で作る。画像で貼らない。
  → 納品後に PowerPoint 上で文言修正・図形サイズ変更ができる状態で渡すため。
  image は写真/スクリーンショット専用（SVGの貼り付けは既定でエラー）。
- フォント最低 16pt（ユーザー必須ルール）。本来11pt/13ptの小見出し・リードも16pt以上に引き上げる。
  figure 内の size も実寸pt。図座標ではないので注意（16未満は自動で16に引き上げ）。
- 主役色ティール #15B5AA は schemeClr 'tx2'（テーマ側で定義済み）を使用。
- 日本語=メイリオ。ロゴ・ページ番号・背景はマスタ/レイアウトが自動で付与。
"""
import sys, json, copy, argparse, os, io, re
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from PIL import Image as PILImage

MIN_PT = 16            # 文字サイズの下限（必須）
TEAL = RGBColor(0x15, 0xB5, 0xAA)
JP_FONT = "メイリオ"

# ---- テンプレ内のスライド位置（template.pptx の固定構成） ----
IDX_COVER   = 0   # 表紙
IDX_TOC     = 1   # 目次
IDX_CONTENT = 2   # 本文グリッド（複製元）
IDX_SECTION = 3   # 章扉（中扉グラフィック・複製元）
IDX_SUMMARY = 4   # サマリーページ（複製元）
IDX_END     = 5   # エンドページ

# 本文スライドの複製元（type → テンプレ位置）
TEMPLATE_INDEX = {"content": IDX_CONTENT, "summary": IDX_SUMMARY}

# 本文／サマリーのコンテンツ領域 (x, y, W, H) EMU（レイアウト定義値）
CONTENT_RECT = (457788, 1707538, 11291300, 4817087)   # 本文グリッド idx12
SUMMARY_RECT = (457787, 2313799, 11303999, 4210825)   # サマリー idx14
COL_GUTTER   = 360000                                  # 2カラム時の間隔

# 章扉バリアント → レイアウト名
SECTION_LAYOUTS = {
    "graphic": "中扉（グラフィック）",
    "photo1": "中扉1（写真）", "photo2": "中扉2（写真）", "photo3": "中扉3（写真）",
    "photo4": "中扉4（写真）", "photo5": "中扉5（写真）", "photo6": "中扉6（写真）",
}

# SVG埋め込み用定数
ASVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
SVG_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"

BASE_DIR = "."          # 画像パス解決の基準（main で spec のあるディレクトリに設定）
ALLOW_SVG_IMAGE = False  # SVGの画像貼り付け（既定禁止。図は figure で作る）


def set_run(run, text, size=None, bold=None, color=None, font=JP_FONT):
    run.text = text
    f = run.font
    if size is not None:
        f.size = Pt(max(size, MIN_PT))
    if bold is not None:
        f.bold = bold
    if color is not None:
        f.color.rgb = color
    if font:
        f.name = font
        # 日本語(ea)フォントも明示
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = rPr.makeelement(qn('a:ea'), {})
            rPr.append(ea)
        ea.set('typeface', font)


def get_ph(slide, ph_type=None, idx=None):
    for sh in slide.placeholders:
        pf = sh.placeholder_format
        if idx is not None and pf.idx == idx:
            return sh
        if ph_type is not None and pf.idx is None and str(pf.type) == ph_type:
            return sh
    return None


def first_textbox(slide):
    """プレースホルダでない最初のテキストボックスを返す（目次の項目ボックス用）。"""
    for sh in slide.shapes:
        if (not sh.is_placeholder) and sh.has_text_frame:
            return sh
    return None


def clear_tf_keep_first(tf):
    """text_frame の段落を最初の1つだけ残して削除（最初をテンプレ段落として使う）。"""
    paras = tf.paragraphs
    for p in paras[1:]:
        p._p.getparent().remove(p._p)
    return tf.paragraphs[0]


def fill_paragraph(p, text, size, bold=False, level=0, bullet=False, color=None):
    """空段落 p にテキストを設定（既存 run を消して1つ作る）。"""
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    p.level = level
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, color=color)
    pPr = p._p.get_or_add_pPr()
    # bullet 制御
    for tag in ('a:buNone', 'a:buChar', 'a:buAutoNum', 'a:buFont'):
        e = pPr.find(qn(tag))
        if e is not None:
            pPr.remove(e)
    if bullet:
        buFont = pPr.makeelement(qn('a:buFont'), {'typeface': 'Arial'})
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '●'})
        pPr.append(buFont); pPr.append(buChar)
        if p.level == 0:
            pPr.set('marL', '228600'); pPr.set('indent', '-228600')
        else:
            pPr.set('marL', str(228600 + 457200 * p.level)); pPr.set('indent', '-228600')
    else:
        pPr.append(pPr.makeelement(qn('a:buNone'), {}))


def duplicate_slide(prs, src_index):
    """プレースホルダのみのスライドを複製（画像 rels なし前提）。末尾に追加して返す。"""
    src = prs.slides[src_index]
    new = prs.slides.add_slide(src.slide_layout)
    # add_slide が作った placeholder を全削除
    for sh in list(new.shapes):
        sh._element.getparent().remove(sh._element)
    # src の図形をコピー
    for sh in src.shapes:
        new.shapes._spTree.append(copy.deepcopy(sh._element))
    return new


def get_layout(prs, name):
    for L in prs.slide_layouts:
        if L.name == name:
            return L
    return None


def resolve_path(p):
    """画像パスを BASE_DIR / cwd / そのまま の順で解決。"""
    if os.path.isabs(p) and os.path.exists(p):
        return p
    for base in (BASE_DIR, "."):
        cand = os.path.join(base, p)
        if os.path.exists(cand):
            return cand
    return p


# ---------------------- 画像 / SVG 配置 ----------------------

def _svg_aspect(svg_bytes):
    s = svg_bytes.decode("utf-8", "ignore")
    m = re.search(r'viewBox\s*=\s*["\']\s*[\d.\-]+\s+[\d.\-]+\s+([\d.]+)\s+([\d.]+)', s)
    if m:
        return float(m.group(1)), float(m.group(2))
    w = re.search(r'\bwidth\s*=\s*["\']([\d.]+)', s)
    h = re.search(r'\bheight\s*=\s*["\']([\d.]+)', s)
    if w and h:
        return float(w.group(1)), float(h.group(1))
    return 16.0, 9.0


def _png_fallback(aw, ah, max_px=1600):
    """簡易PNGフォールバック（最新PowerPointは本物のSVGを描画するので保険）。"""
    scale = max_px / max(aw, ah)
    w, h = max(1, int(aw * scale)), max(1, int(ah * scale))
    buf = io.BytesIO()
    PILImage.new("RGBA", (w, h), (255, 255, 255, 0)).save(buf, "PNG")
    buf.seek(0)
    return buf


def _fit(rect, aw, ah):
    """rect=(x,y,W,H) にアスペクト比を保って内接させ中央寄せ。"""
    x, y, W, H = rect
    s = min(W / aw, H / ah)
    w, h = int(aw * s), int(ah * s)
    return x + (W - w) // 2, y + (H - h) // 2, w, h


def _add_svg(slide, svg_path, rect):
    with open(svg_path, "rb") as f:
        svg_bytes = f.read()
    aw, ah = _svg_aspect(svg_bytes)
    x, y, w, h = _fit(rect, aw, ah)
    pic = slide.shapes.add_picture(_png_fallback(aw, ah), Emu(x), Emu(y), Emu(w), Emu(h))
    pkg = slide.part.package
    partname = pkg.next_partname("/ppt/media/image%d.svg")
    svg_part = Part(partname, "image/svg+xml", pkg, svg_bytes)
    rId = slide.part.relate_to(svg_part, RT.IMAGE)
    blip = pic._element.find(".//" + qn("a:blip"))
    extLst = blip.makeelement(qn("a:extLst"), {})
    ext = blip.makeelement(qn("a:ext"), {"uri": SVG_EXT_URI})
    svgBlip = ext.makeelement("{%s}svgBlip" % ASVG_NS, {})
    svgBlip.set(qn("r:embed"), rId)
    ext.append(svgBlip)
    extLst.append(ext)
    blip.append(extLst)
    return pic


def add_image(slide, img_path, rect):
    """写真/スクリーンショットをアスペクト維持で rect に配置。

    図解を画像で貼ると納品後に文言もサイズも直せなくなるため、SVGの貼り付けは既定で禁止。
    図解は figure（ネイティブ図形）で作ること。
    """
    path = resolve_path(img_path)
    if not os.path.exists(path):
        return None
    if os.path.splitext(path)[1].lower() == ".svg":
        if not ALLOW_SVG_IMAGE:
            raise SystemExit(
                f"エラー: SVGの画像貼り付けは禁止されています → {img_path}\n"
                "  図解は image ではなく figure（ネイティブ図形）で作ってください。\n"
                "  画像で貼ると、PowerPoint上で文言修正も図形サイズ変更もできなくなります。\n"
                "  どうしても画像で貼る場合のみ --allow-svg-image を付けてください。")
        return _add_svg(slide, path, rect)
    with PILImage.open(path) as im:
        aw, ah = im.size
    x, y, w, h = _fit(rect, aw, ah)
    return slide.shapes.add_picture(path, Emu(x), Emu(y), Emu(w), Emu(h))


def set_box(ph, rect):
    """プレースホルダの位置・サイズを rect=(x,y,W,H) に固定。"""
    x, y, w, h = rect
    ph.left, ph.top, ph.width, ph.height = Emu(x), Emu(y), Emu(w), Emu(h)


# ------------------- 図（PowerPointネイティブ図形） -------------------
# 図は画像で貼らずに、必ずオートシェイプ＋テキストフレームで作る。
# こうしておくと納品後に PowerPoint 上で文言修正・サイズ変更・色替えができる。

FIG_COLORS = {
    "teal": "15B5AA", "tealDark": "0F8F87", "pale": "EFF5F5", "paleLine": "C9DCDC",
    "white": "FFFFFF", "black": "000000", "ink": "5D6B78",
    "gray": "F0F0F0", "grayLine": "BBBBBB", "dark": "222222",
    "blue": "7392F5", "cyan": "88DBDF", "purple": "8557A7",
    "green": "96DC00", "red": "FF5A41", "yellow": "FFB41E",
}
# 濃い塗り＝文字を白に自動反転する塗り
DARK_FILLS = {"teal", "tealDark", "dark", "black", "purple", "blue", "red"}

FIG_SHAPES = {
    "rect": MSO_SHAPE.RECTANGLE,
    "roundRect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "chevron": MSO_SHAPE.CHEVRON,
    "arrowBox": MSO_SHAPE.RIGHT_ARROW,
    "pentagon": MSO_SHAPE.PENTAGON,
    "diamond": MSO_SHAPE.DIAMOND,
}

ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
VALIGN = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}


def fig_color(name, default=None):
    """色名 or #RRGGBB を RGBColor に。'none'/None は None（＝塗り/線なし）。"""
    if name is None:
        return default
    if not isinstance(name, str) or name.lower() == "none":
        return None
    if name.startswith("#"):
        return RGBColor.from_string(name[1:])
    if name in FIG_COLORS:
        return RGBColor.from_string(FIG_COLORS[name])
    return default


def _fig_text(shape, item, default_size=16, default_align="c", default_valign="m"):
    """図形/テキストボックスに文字を流し込む（複数行は lines か \\n）。"""
    lines = item.get("lines")
    if lines is None:
        raw = item.get("text")
        lines = str(raw).split("\n") if raw not in (None, "") else []
    if not lines:
        return
    tf = shape.text_frame
    tf.word_wrap = item.get("wrap", True)
    tf.margin_left = tf.margin_right = Emu(45720)   # 0.05in
    tf.margin_top = tf.margin_bottom = Emu(18288)
    tf.vertical_anchor = VALIGN.get(item.get("valign", default_valign), MSO_ANCHOR.MIDDLE)

    # 塗りが濃ければ文字は白（color 明示があればそちら優先）
    fill_name = item.get("fill")
    auto = "white" if (isinstance(fill_name, str) and fill_name in DARK_FILLS) else "black"
    color = fig_color(item.get("color"), fig_color(auto))
    size = item.get("size", default_size)
    align = ALIGN.get(item.get("align", default_align), PP_ALIGN.CENTER)

    for i, ln in enumerate(lines):
        if isinstance(ln, str):
            ln = {"text": ln}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        p.alignment = ALIGN.get(ln.get("align"), align)
        pPr = p._p.get_or_add_pPr()
        pPr.append(pPr.makeelement(qn('a:buNone'), {}))
        set_run(p.add_run(), ln.get("text", ""),
                size=ln.get("size", size),
                bold=ln.get("bold", item.get("bold", False)),
                color=fig_color(ln.get("color"), color))


def _fig_box(slide, item, tx):
    """矩形・角丸・楕円などのオートシェイプを1つ作る。"""
    kind = item.get("kind", "roundRect")
    x, y, w, h = tx(item.get("x", 0), item.get("y", 0), item.get("w", 100), item.get("h", 40))
    if kind == "text":
        shape = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
        _fig_text(shape, item, default_align=item.get("align", "l"),
                  default_valign=item.get("valign", "t"))
        return shape

    shape = slide.shapes.add_shape(FIG_SHAPES.get(kind, MSO_SHAPE.ROUNDED_RECTANGLE),
                                   Emu(x), Emu(y), Emu(w), Emu(h))
    shape.shadow.inherit = False        # テンプレの影を継がない

    fill = fig_color(item.get("fill", "white"))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill

    border = fig_color(item.get("line"))
    if border is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border
        shape.line.width = Pt(item.get("lineWidth", 1.5))

    # 角丸の丸み: radius(図座標) / 短辺
    if kind == "roundRect" and "radius" in item:
        base = max(1, min(item.get("w", 100), item.get("h", 40)))
        try:
            shape.adjustments[0] = max(0.0, min(0.5, item["radius"] / base))
        except (IndexError, KeyError):
            pass

    _fig_text(shape, item)
    return shape


def _fig_line(slide, item, tx):
    """折れ線・矢印。points=[[x,y],...] を直線コネクタの連なりで描く。"""
    pts = item.get("points")
    if not pts or len(pts) < 2:
        return None
    color = fig_color(item.get("color"), fig_color("teal"))
    width = Pt(item.get("width", 2))
    arrow = item.get("arrow", item.get("kind") == "arrow")
    out = []
    for i in range(len(pts) - 1):
        (x1, y1), (x2, y2) = pts[i], pts[i + 1]
        ax, ay, _, _ = tx(x1, y1, 0, 0)
        bx, by, _, _ = tx(x2, y2, 0, 0)
        cxn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                         Emu(ax), Emu(ay), Emu(bx), Emu(by))
        cxn.line.color.rgb = color
        cxn.line.width = width
        ln = cxn.line._get_or_add_ln()
        if item.get("dash"):
            ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
        # 最後の線分にだけ矢じりを付ける（a:ln の子要素順序に従い末尾へ）
        if arrow and i == len(pts) - 2:
            ln.append(ln.makeelement(qn('a:tailEnd'),
                                     {'type': 'triangle', 'w': 'med', 'len': 'med'}))
        out.append(cxn)
    return out


def add_figure(slide, fig, rect):
    """図座標(viewBox)を rect=(x,y,W,H) にアスペクト維持で内接させて図形を配置。

    fig = {"viewBox":[w,h] または [x,y,w,h], "shapes":[...]}
    テキストサイズだけは図座標ではなく実寸pt（16pt以上）で指定する。
    """
    if not isinstance(fig, dict):
        return
    vb = fig.get("viewBox", [1180, 500])
    vx, vy, vw, vh = (0, 0, vb[0], vb[1]) if len(vb) == 2 else vb
    if not vw or not vh:
        return
    ox, oy, w, h = _fit(rect, vw, vh)
    s = w / float(vw)

    def tx(x, y, bw, bh):
        return (int(ox + (x - vx) * s), int(oy + (y - vy) * s),
                max(1, int(bw * s)), max(1, int(bh * s)))

    for item in fig.get("shapes", []):
        if not isinstance(item, dict):
            continue
        kind = item.get("kind", "roundRect")
        if kind in ("line", "arrow") or "points" in item:
            _fig_line(slide, item, tx)
        else:
            _fig_box(slide, item, tx)


def delete_slide(prs, sldId_el):
    """sldId要素・リレーション・パートをまとめて削除（修復警告を防ぐ）。"""
    rId = sldId_el.get(qn('r:id'))
    if rId:
        try:
            prs.part.drop_rel(rId)
        except KeyError:
            pass
    prs.slides._sldIdLst.remove(sldId_el)


# ----------------------- 各スライド生成 -----------------------

def build_cover(slide, spec):
    t = get_ph(slide, idx=0)
    if t:
        p = t.text_frame.paragraphs[0]
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        set_run(p.add_run(), spec.get("title", "プレゼンテーションタイトル"), bold=True)
    sub = get_ph(slide, idx=1)
    if sub:
        p = sub.text_frame.paragraphs[0]
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        set_run(p.add_run(), spec.get("subtitle", ""), size=18)
    date = get_ph(slide, idx=13)
    if date:
        p = date.text_frame.paragraphs[0]
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        set_run(p.add_run(), spec.get("date", ""), size=16)


def fill_body(slide, idx, lines):
    """本文プレースホルダ(idx)に行リストを流し込む。"""
    body = get_ph(slide, idx=idx)
    if not body:
        return
    tf = body.text_frame
    first = clear_tf_keep_first(tf)
    if not lines:
        fill_paragraph(first, "", size=18)
        return
    for i, ln in enumerate(lines):
        if isinstance(ln, str):
            ln = {"text": ln}
        p = first if i == 0 else tf.add_paragraph()
        fill_paragraph(p, ln.get("text", ""),
                       size=ln.get("size", 18),
                       bold=ln.get("bold", False),
                       level=ln.get("level", 0),
                       bullet=ln.get("bullet", False))


def fill_content_area(slide, body_idx, rect, item):
    """コンテンツ領域を埋める。
    - 図/画像のみ     → 領域いっぱいに配置
    - 本文＋図/画像   → 左=本文 / 右=図 の2カラム
    - 本文のみ        → 従来通り

    図は figure（ネイティブ図形・編集可）を優先。image は写真/スクショ専用。
    """
    figure = item.get("figure")
    image = item.get("image")
    body = item.get("body", [])
    body_ph = get_ph(slide, idx=body_idx)
    x, y, W, H = rect

    def place(r):
        if figure:
            add_figure(slide, figure, r)
        elif image:
            add_image(slide, image, r)

    if (figure or image) and body:
        body_w = int(W * 0.50)
        if body_ph:
            set_box(body_ph, (x, y, body_w, H))
        fill_body(slide, body_idx, body)
        place((x + body_w + COL_GUTTER, y, W - body_w - COL_GUTTER, H))
    elif figure or image:
        if body_ph:                                   # 空の本文枠は黙らせる
            clear_tf_keep_first(body_ph.text_frame)
            fill_paragraph(body_ph.text_frame.paragraphs[0], "", size=18)
        place(rect)
    else:
        fill_body(slide, body_idx, body)


def build_content(slide, item):
    sec = get_ph(slide, idx=10)     # セクション小見出し
    if sec:
        fill_paragraph(sec.text_frame.paragraphs[0], item.get("section", ""),
                       size=16, bold=True, color=TEAL)
    title = get_ph(slide, idx=0)    # ページタイトル
    if title:
        fill_paragraph(title.text_frame.paragraphs[0], item.get("title", ""),
                       size=26, bold=True)
    lead = get_ph(slide, idx=11)    # リード文
    if lead:
        fill_paragraph(lead.text_frame.paragraphs[0], item.get("lead", ""), size=16)
    fill_content_area(slide, 12, CONTENT_RECT, item)


def build_summary(slide, item):
    """サマリーページ（背景画像＋タイトル＋リード＋本文 idx=14）。"""
    sec = get_ph(slide, idx=10)
    if sec:
        fill_paragraph(sec.text_frame.paragraphs[0], item.get("section", ""),
                       size=16, bold=True, color=TEAL)
    title = get_ph(slide, idx=0)
    if title:
        fill_paragraph(title.text_frame.paragraphs[0], item.get("title", ""),
                       size=26, bold=True)
    lead = get_ph(slide, idx=11)
    if lead:
        fill_paragraph(lead.text_frame.paragraphs[0], item.get("lead", ""), size=16)
    fill_content_area(slide, 14, SUMMARY_RECT, item)


def build_section(slide, item, auto_number):
    """章扉（中扉グラフィック）。大きな番号(idx=11, 約100pt ティール)＋タイトル(idx=0, 36pt太字)。"""
    title = get_ph(slide, idx=0)
    if title:
        fill_paragraph(title.text_frame.paragraphs[0], item.get("title", ""),
                       size=36, bold=True)
    num = get_ph(slide, idx=11)
    if num:
        # サイズ・色はレイアウト(100pt/ティール)を継承させるため size=None
        fill_paragraph(num.text_frame.paragraphs[0],
                       str(item.get("number", auto_number)), size=None)


def build_toc(slide, entries):
    """entries = [(title, page), ...]。page は表示ページ番号(0始まり)。"""
    box = first_textbox(slide)
    if box is None:
        return
    tf = box.text_frame
    template_p = tf.paragraphs[0]._p           # 既存の点線リーダー付き項目をテンプレに
    for p in list(tf.paragraphs):
        p._p.getparent().remove(p._p)
    for n, (title, page) in enumerate(entries, 1):
        newp = copy.deepcopy(template_p)
        tf._txBody.append(newp)
        # タブ停止を右端1つに統一（短いタイトルでもページ番号を右揃えに）
        pPr = newp.find(qn('a:pPr'))
        if pPr is not None:
            tabLst = pPr.find(qn('a:tabLst'))
            if tabLst is not None:
                for t in list(tabLst):
                    tabLst.remove(t)
                tabLst.append(tabLst.makeelement(qn('a:tab'), {'pos': '7285038', 'algn': 'l'}))
        runs = newp.findall(qn('a:r'))
        if not runs:
            continue
        # run 構成: [番号, 　タイトル, (\tリーダー), ...ページ桁]
        leader_i = None
        for i, r in enumerate(runs):
            t = r.find(qn('a:t'))
            if t is not None and t.text and '\t' in t.text:
                leader_i = i
                break
        t0 = runs[0].find(qn('a:t'))
        if t0 is not None:
            t0.text = f"{n:02d}"
        if len(runs) > 1:
            t1 = runs[1].find(qn('a:t'))
            if t1 is not None:
                t1.text = "　" + title
        if leader_i is not None and leader_i + 1 < len(runs):
            after = runs[leader_i + 1:]
            tafter = after[0].find(qn('a:t'))
            if tafter is not None:
                tafter.text = f" {page:02d}"
            for extra in after[1:]:
                extra.getparent().remove(extra)


def build_end(slide, spec):
    if "thanks" in spec:
        sub = get_ph(slide, idx=1)
        if sub:
            p = sub.text_frame.paragraphs[0]
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
            set_run(p.add_run(), spec["thanks"], size=16)


# ページ番号・日付・フッターは中身が空でも消さない（レイアウト由来の機能枠）
KEEP_PH_TYPES = {"SLIDE_NUMBER", "DATE", "FOOTER"}


def drop_empty_placeholders(prs):
    """使わなかったプレースホルダを削除する。

    空のまま残すと編集画面に「本文/図表/グラフなど」等のプロンプト文字が表示され、
    受け取った人の邪魔になる（書き出し画像には出ないので目視QAでは気づけない）。
    """
    removed = 0
    for s in prs.slides:
        for sh in list(s.shapes):
            if not sh.is_placeholder or not sh.has_text_frame:
                continue                                  # 画像プレースホルダ等は対象外
            if str(sh.placeholder_format.type).split(" ")[0] in KEEP_PH_TYPES:
                continue
            if sh.text_frame.text.strip():
                continue                                  # 中身があるものは残す
            sh._element.getparent().remove(sh._element)
            removed += 1
    return removed


def enforce_min_font(prs):
    """全 run を走査し、明示サイズが MIN_PT 未満なら引き上げる（保険）。"""
    for s in prs.slides:
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size is not None and r.font.size < Pt(MIN_PT):
                        r.font.size = Pt(MIN_PT)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("spec")
    ap.add_argument("output")
    ap.add_argument("--template", default=None)
    ap.add_argument("--allow-svg-image", action="store_true",
                    help="SVGを画像として貼る（既定は禁止。図は figure で作ること）")
    args = ap.parse_args()

    global ALLOW_SVG_IMAGE
    ALLOW_SVG_IMAGE = args.allow_svg_image

    here = os.path.dirname(os.path.abspath(__file__))
    template = args.template or os.path.join(here, "template.pptx")

    with open(args.spec, encoding="utf-8-sig") as f:   # BOM付きでも許容
        spec = json.load(f)

    global BASE_DIR
    BASE_DIR = os.path.dirname(os.path.abspath(args.spec))   # 画像相対パスの基準

    prs = Presentation(template)
    items = spec.get("slides", [])
    if not items:
        raise SystemExit("spec.slides が空です。本文スライドを1枚以上指定してください。")

    use_toc = spec.get("toc", len(items) >= 2)

    # 表紙・エンド（テンプレの固定スライドを編集）
    build_cover(prs.slides[IDX_COVER], spec)
    build_end(prs.slides[IDX_END], spec)

    sldIdLst = prs.slides._sldIdLst
    orig = list(sldIdLst)   # [表紙, 目次, 本文雛形, 章扉雛形, サマリー雛形, エンド]
    cover_id, toc_id = orig[0], orig[1]
    tmpl_ids = orig[IDX_CONTENT:IDX_SUMMARY + 1]   # 複製元3枚（削除対象）
    end_id = orig[IDX_END]

    # 各スライドを type に応じて生成（末尾に追加される）
    out = []
    auto_sec = 0
    for item in items:
        typ = item.get("type", "content")
        if typ == "section":
            name = SECTION_LAYOUTS.get(item.get("variant", "graphic"), SECTION_LAYOUTS["graphic"])
            layout = get_layout(prs, name) or prs.slides[IDX_SECTION].slide_layout
            ns = prs.slides.add_slide(layout)   # 章扉は選んだレイアウトで新規作成
            new_id = list(sldIdLst)[-1]
            auto_sec += 1
            build_section(ns, item, f"{auto_sec:02d}")
        elif typ == "summary":
            ns = duplicate_slide(prs, IDX_SUMMARY)
            new_id = list(sldIdLst)[-1]
            build_summary(ns, item)
        else:
            ns = duplicate_slide(prs, IDX_CONTENT)
            new_id = list(sldIdLst)[-1]
            build_content(ns, item)
        out.append({"id": new_id, "type": typ, "title": item.get("title", "")})

    # 最終順序を確定し、表示ページ番号(0始まり)を算出
    final = [cover_id] + ([toc_id] if use_toc else []) + [o["id"] for o in out] + [end_id]
    page_of = {el: pos for pos, el in enumerate(final)}

    # 目次（章扉があれば章扉を、無ければ本文/サマリーのタイトルを列挙）
    if use_toc:
        has_section = any(o["type"] == "section" for o in out)
        listed = [o for o in out if (o["type"] == "section") == has_section]
        build_toc(prs.slides[IDX_TOC], [(o["title"], page_of[o["id"]]) for o in listed])

    # 複製元の雛形3枚を削除、目次不要なら削除
    for tid in tmpl_ids:
        delete_slide(prs, tid)
    if not use_toc:
        delete_slide(prs, toc_id)

    # 並べ替え
    for el in list(sldIdLst):
        sldIdLst.remove(el)
    for el in final:
        sldIdLst.append(el)

    dropped = drop_empty_placeholders(prs)
    enforce_min_font(prs)
    prs.save(args.output)
    print(f"OK: {args.output}  slides={len(list(prs.slides._sldIdLst))}  "
          f"空プレースホルダ削除={dropped}")


if __name__ == "__main__":
    main()
