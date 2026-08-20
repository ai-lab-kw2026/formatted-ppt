#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
既存スライド（.pptx / .pdf）から中身を機械的に抽出し、digest.json を書き出す。

このスクリプトは「判断しない」。
住商ルールでの再構成（言い切りタイトルへの書き換え・1枚1メッセージへの分割・
図の figure 化）は SKILL.md の「変換モード」に従って Claude が行う。
ここでやるのは、そのための材料を漏れなく・忠実に取り出すことだけ。

    python extract.py deck.pptx -o digest.json [--media-dir media] [--md digest.md]
    python extract.py deck.pdf  -o digest.json

出力 digest.json:
    {
      "source": "deck.pptx", "kind": "pptx",
      "slideSize": {"w": 12192000, "h": 6858000, "aspect": 1.777},
      "slides": [
        {"index": 1, "layout": "タイトルとコンテンツ", "notes": "...",
         "shapes": [
           {"type": "text", "roleHint": "title", "pos": {...},
            "paragraphs": [{"text": "...", "level": 0, "bullet": false,
                            "size": 40.0, "bold": true}]},
           {"type": "table", "roleHint": "body", "pos": {...},
            "rows": [["列1", "列2"], ["a", "b"]]},
           {"type": "picture", "roleHint": "body", "pos": {...},
            "file": "media/s01_p1.png", "w": 800, "h": 600,
            "aspect": 1.33, "hint": "photo"},
           {"type": "shape", "roleHint": "decoration", "pos": {...},
            "shapeKind": "roundRect", "fill": "teal", "text": "..."}
         ]}
      ]
    }

pos は EMU をスライド比 0〜1 に正規化した {x, y, w, h}。
テンプレートが違っても意味が保てるようにするため。
"""

import argparse
import json
import os
import sys

try:
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from pptx.util import Emu
except ImportError:
    raise SystemExit(
        "エラー: python-pptx が見つかりません。\n"
        "  python -m pip install python-pptx pillow"
    )

# build.py のブランド定義を再利用する（色名・図形名の逆引きに使う）。
# build.py は if __name__ == "__main__" ガード付きなので import しても実行されない。
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
try:
    from build import FIG_COLORS, FIG_SHAPES
except Exception:  # build.py が隣に無い/壊れている場合も抽出自体は続行させる
    FIG_COLORS = {}
    FIG_SHAPES = {}

# hex → ブランド色名（"15B5AA" → "teal"）
HEX_TO_NAME = {v.upper(): k for k, v in FIG_COLORS.items()}
# MSO_SHAPE → figure の kind（ROUNDED_RECTANGLE → "roundRect"）
MSO_TO_KIND = {}
for _k, _v in FIG_SHAPES.items():
    try:
        MSO_TO_KIND[int(_v)] = _k
    except (TypeError, ValueError):
        pass

# 図形として描き直す候補になりやすい画像形式（ベクタ・図版）
VECTOR_EXT = {".emf", ".wmf", ".svg"}
# 写真として引き継ぐ候補になりやすい形式
PHOTO_EXT = {".jpg", ".jpeg"}

# GraphicFrame の中身を見分けるための名前空間
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_DGM = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
NS_OLE = "http://schemas.openxmlformats.org/presentationml/2006/ole"
# graphicData/@uri → 人が読める名前
GRAPHIC_KINDS = {
    NS_DGM: "SmartArt",
    NS_OLE: "OLE埋め込みオブジェクト",
    "http://schemas.openxmlformats.org/drawingml/2006/table": "表",
    "http://schemas.openxmlformats.org/drawingml/2006/chart": "グラフ",
}


# ---------------------------------------------------------------- 位置・書式

def _norm_pos(shape, sw, sh):
    """EMU の位置・サイズをスライド比 0〜1 に正規化する。"""

    def _r(v, base):
        if v is None or not base:
            return None
        return round(float(v) / float(base), 4)

    return {
        "x": _r(shape.left, sw),
        "y": _r(shape.top, sh),
        "w": _r(shape.width, sw),
        "h": _r(shape.height, sh),
    }


def _bullet_of(para):
    """段落が箇条書きかを a:pPr の bullet 要素から判定する。

    build.py が書き出す側で使う a:buChar / marL を、読む側から見た逆引き。
    """
    pPr = para._p.find(qn("a:pPr"))
    if pPr is None:
        return False
    if pPr.find(qn("a:buNone")) is not None:
        return False
    for tag in ("a:buChar", "a:buAutoNum"):
        if pPr.find(qn(tag)) is not None:
            return True
    return False


def _pt(size):
    """python-pptx の Length を pt の float にする。None はそのまま。"""
    if size is None:
        return None
    try:
        return round(size.pt, 1)
    except AttributeError:
        return None


def _para_format(para):
    """段落の代表的な文字サイズ・太字を、最初の実体のある run から拾う。"""
    size = _pt(para.font.size)
    bold = para.font.bold
    for run in para.runs:
        if not run.text.strip():
            continue
        if size is None:
            size = _pt(run.font.size)
        if bold is None:
            bold = run.font.bold
        if size is not None and bold is not None:
            break
    return size, bool(bold) if bold is not None else False


def _read_paragraphs(tf):
    """テキストフレームを段落の配列にする。空段落は落とす。"""
    out = []
    for para in tf.paragraphs:
        text = "".join(r.text for r in para.runs) or para.text
        text = text.strip()
        if not text:
            continue
        size, bold = _para_format(para)
        out.append({
            "text": text,
            "level": int(para.level or 0),
            "bullet": _bullet_of(para),
            "size": size,
            "bold": bold,
        })
    return out


def _fill_name(shape):
    """図形の塗りを、可能ならブランド色名に、無理なら #RRGGBB で返す。"""
    try:
        fill = shape.fill
        if fill.type is None or int(fill.type) != 1:  # 1 = MSO_FILL.SOLID
            return None
        rgb = fill.fore_color.rgb
    except Exception:
        return None
    if rgb is None:
        return None
    hexv = str(rgb).upper()
    return HEX_TO_NAME.get(hexv, "#" + hexv)


def _shape_kind(shape):
    """オートシェイプの種類を figure の kind 名に寄せる。不明なら None。"""
    try:
        auto = shape.auto_shape_type
    except Exception:
        return None
    if auto is None:
        return None
    return MSO_TO_KIND.get(int(auto))


# ---------------------------------------------------------------- 役割の当たり

def _role_hint(entry, max_size):
    """位置と文字サイズから役割の当たりをつける。

    確定ではなく Claude への手がかり。誤ってもよいが、上から順に読めば
    元スライドの意図が追える程度には当てる。
    """
    if entry.get("isTitlePlaceholder"):
        return "title"

    pos = entry.get("pos") or {}
    y = pos.get("y")
    h = pos.get("h") or 0
    size = None
    for para in entry.get("paragraphs") or []:
        if para.get("size") is not None:
            size = max(size or 0, para["size"])

    # 上部にあって、そのスライドで最も大きい文字 → タイトル
    if y is not None and y < 0.28 and size is not None and max_size and size >= max_size:
        return "title"
    # 下部の小さい文字 → 出典・注釈
    if y is not None and y > 0.82 and (size is None or size <= 12):
        return "caption"
    # 文字が無く、面積も小さい図形 → 飾り
    if entry["type"] == "shape" and not entry.get("text"):
        if (pos.get("w") or 0) * (h or 0) < 0.02:
            return "decoration"
    return "body"


# ---------------------------------------------------------------- 画像

def _picture_hint(path):
    """写真か、figure で描き直すべき図版かの当たりをつける。

    ベクタ形式と、色数の少ない PNG は図版とみなす。
    """
    ext = os.path.splitext(path)[1].lower()
    if ext in VECTOR_EXT:
        return "figure_candidate"
    if ext in PHOTO_EXT:
        return "photo"
    try:
        from PIL import Image
    except ImportError:
        return "unknown"
    try:
        with Image.open(path) as im:
            im = im.convert("RGB")
            im.thumbnail((160, 160))
            colors = im.getcolors(maxcolors=4096)
    except Exception:
        return "unknown"
    if colors is None:          # 4096 色を超える = 写真
        return "photo"
    return "figure_candidate" if len(colors) <= 64 else "photo"


def _save_picture(shape, media_dir, slide_no, seq):
    """画像を media_dir に書き出し、(相対パス, 幅, 高さ) を返す。"""
    try:
        image = shape.image
    except Exception:
        return None, None, None
    ext = image.ext or "png"
    name = "s%02d_p%d.%s" % (slide_no, seq, ext)
    os.makedirs(media_dir, exist_ok=True)
    path = os.path.join(media_dir, name)
    with open(path, "wb") as fh:
        fh.write(image.blob)
    try:
        px = image.size
    except Exception:
        px = (None, None)
    return path, px[0], px[1]


# ------------------------------------------------- SmartArt / OLE など特殊枠

def _graphic_uri(shape):
    """GraphicFrame の graphicData/@uri を返す。GraphicFrame でなければ None。"""
    el = shape.element
    if el.tag != "{http://schemas.openxmlformats.org/presentationml/2006/main}graphicFrame":
        return None
    gd = el.find(".//{%s}graphicData" % NS_A)
    return gd.get("uri") if gd is not None else None


def _smartart_nodes(shape):
    """SmartArt のテキストを取り出す。

    SmartArt の文字はスライド本体ではなく ppt/diagrams/data#.xml にある。
    graphicFrame の dgm:relIds/@r:dm がその関連付けを指しているので、
    リレーションを辿ってノードごとのテキストを拾う。
    """
    rel_ids = shape.element.find(".//{%s}relIds" % NS_DGM)
    if rel_ids is None:
        return []
    rid = rel_ids.get("{%s}dm" % NS_R)
    if not rid:
        return []
    try:
        blob = shape.part.rels[rid].target_part.blob
    except Exception:
        return []
    try:
        from lxml import etree
        root = etree.fromstring(blob)
    except Exception:
        return []

    nodes, seen = [], set()
    for pt in root.iter("{%s}pt" % NS_DGM):
        lines = []
        for para in pt.iter("{%s}p" % NS_A):
            text = "".join(t.text or "" for t in para.iter("{%s}t" % NS_A)).strip()
            if text:
                lines.append(text)
        if not lines:
            continue
        joined = "\n".join(lines)
        # 同じ内容の点（プレゼン用の複製）が並ぶことがあるので重複は落とす
        if joined not in seen:
            seen.add(joined)
            nodes.append(joined)
    return nodes


def _graphic_entry(shape, uri, pos):
    """表・グラフ以外の GraphicFrame を digest の1項目にする。

    ここで拾えなくても「拾えなかった」ことは必ず残す。無言で落とすと、
    元スライドに図があった事実そのものが変換時に見えなくなるため。
    """
    kind = GRAPHIC_KINDS.get(uri, "不明な埋め込み")

    if uri == NS_DGM:
        nodes = _smartart_nodes(shape)
        entry = {"type": "smartart", "pos": pos, "nodes": nodes}
        if nodes:
            entry["text"] = "\n".join(nodes)
            entry["note"] = ("SmartArt。テキストは取得済み。figure で描き直すこと"
                             "（元の図形は引き継がれない）。")
        else:
            entry["note"] = ("SmartArt だがテキストを取得できなかった。"
                             "元ファイルを開いて内容を確認し、figure で描き直すこと。")
        return entry

    return {
        "type": "unsupported",
        "pos": pos,
        "graphic": kind,
        "graphicUri": uri,
        "note": ("%s は抽出できない。元ファイルを開いて内容を確認し、"
                 "figure で描き直すか、写真なら画像として貼ること。" % kind),
    }


# ---------------------------------------------------------------- 図形の走査

def _walk(shapes, sw, sh, slide_no, media_dir, out, counter):
    """図形ツリーを z 順に走査する。グループは中身に展開する。"""
    for shape in shapes:
        # shape_type は種類によっては例外を投げるので、一度だけ安全に取る
        try:
            stype = str(shape.shape_type)
        except Exception:
            stype = ""

        # グループは再帰。グループ自体は記録しない（中身が本体なので）
        if "GROUP" in stype:
            try:
                _walk(shape.shapes, sw, sh, slide_no, media_dir, out, counter)
            except Exception:
                pass
            continue

        pos = _norm_pos(shape, sw, sh)

        # 表
        if getattr(shape, "has_table", False):
            rows = []
            for row in shape.table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            out.append({"type": "table", "pos": pos, "rows": rows})
            continue

        # グラフ（数値だけ拾う。figure/表で描き直す材料）
        if getattr(shape, "has_chart", False):
            chart = {"type": "chart", "pos": pos, "series": []}
            try:
                plot = shape.chart.plots[0]
                chart["categories"] = [str(c) for c in plot.categories]
                for ser in plot.series:
                    chart["series"].append(
                        {"name": ser.name, "values": list(ser.values)})
            except Exception:
                pass
            out.append(chart)
            continue

        # 表・グラフ以外の GraphicFrame（SmartArt / OLE など）
        uri = _graphic_uri(shape)
        if uri is not None:
            out.append(_graphic_entry(shape, uri, pos))
            continue

        # 画像
        if "PICTURE" in stype:
            counter["pic"] += 1
            if media_dir:
                path, pw, ph = _save_picture(
                    shape, media_dir, slide_no, counter["pic"])
            else:
                path, pw, ph = None, None, None
            entry = {"type": "picture", "pos": pos, "file": path,
                     "w": pw, "h": ph}
            if pw and ph:
                entry["aspect"] = round(pw / float(ph), 3)
            entry["hint"] = _picture_hint(path) if path else "unknown"
            out.append(entry)
            continue

        # テキスト / オートシェイプ
        has_tf = getattr(shape, "has_text_frame", False)
        paras = _read_paragraphs(shape.text_frame) if has_tf else []
        kind = _shape_kind(shape)
        fill = _fill_name(shape)

        if not paras and kind is None and fill is None:
            continue  # 中身も形も色も無いものは記録しない

        is_title = False
        try:
            if shape.is_placeholder:
                ph_type = str(shape.placeholder_format.type or "")
                is_title = "TITLE" in ph_type
        except Exception:
            pass

        entry = {
            "type": "shape" if kind or fill else "text",
            "pos": pos,
            "paragraphs": paras,
        }
        if paras:
            entry["text"] = "\n".join(p["text"] for p in paras)
        if kind:
            entry["shapeKind"] = kind
        if fill:
            entry["fill"] = fill
        if is_title:
            entry["isTitlePlaceholder"] = True
        out.append(entry)


# ---------------------------------------------------------------- pptx / pdf

def extract_pptx(path, media_dir):
    prs = Presentation(path)
    sw, sh = prs.slide_width, prs.slide_height
    digest = {
        "source": os.path.basename(path),
        "kind": "pptx",
        "slideSize": {
            "w": int(sw), "h": int(sh),
            "aspect": round(float(sw) / float(sh), 3) if sh else None,
        },
        "slides": [],
    }

    for i, slide in enumerate(prs.slides, start=1):
        shapes = []
        counter = {"pic": 0}
        _walk(slide.shapes, sw, sh, i, media_dir, shapes, counter)

        # そのスライドの最大文字サイズを出してから役割を当てる
        max_size = 0
        for entry in shapes:
            for para in entry.get("paragraphs") or []:
                if para.get("size"):
                    max_size = max(max_size, para["size"])
        for entry in shapes:
            entry["roleHint"] = _role_hint(entry, max_size)

        # タイトルは1枚に1つだけ。複数当たったら最も上のものを残す
        titles = [e for e in shapes if e["roleHint"] == "title"]
        if len(titles) > 1:
            titles.sort(key=lambda e: (e.get("pos") or {}).get("y") or 0)
            for entry in titles[1:]:
                entry["roleHint"] = "subtitle"

        # 読み順に並べ替える（上から下、同じ高さなら左から右）
        shapes.sort(key=lambda e: (
            round((e.get("pos") or {}).get("y") or 0, 2),
            (e.get("pos") or {}).get("x") or 0,
        ))

        item = {
            "index": i,
            "layout": slide.slide_layout.name,
            "shapes": shapes,
        }
        try:
            if slide.has_notes_slide:
                note = slide.notes_slide.notes_text_frame.text.strip()
                if note:
                    item["notes"] = note
        except Exception:
            pass
        digest["slides"].append(item)

    return digest


def extract_pdf(path):
    try:
        from pypdf import PdfReader
    except ImportError:
        raise SystemExit(
            "エラー: PDFの読み取りには pypdf が必要です → %s\n"
            "  python -m pip install pypdf\n"
            "  ※ .pptx の変換だけなら pypdf は不要です。" % path
        )
    except BaseException as exc:
        # pypdf は入っているが、依存(cryptography/cffi)が壊れていて import が落ちる場合。
        # ImportError では捕まらない種類の例外が飛んでくるのでここで受ける。
        raise SystemExit(
            "エラー: pypdf の読み込みに失敗しました → %s\n"
            "  依存パッケージが壊れている可能性があります。次を試してください:\n"
            "    python -m pip install --upgrade --force-reinstall cffi cryptography pypdf\n"
            "  詳細: %s: %s" % (path, type(exc).__name__, exc)
        )

    reader = PdfReader(path)
    digest = {
        "source": os.path.basename(path),
        "kind": "pdf",
        "slideSize": None,
        "slides": [],
    }
    for i, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
        shapes = []
        if lines:
            shapes.append({
                "type": "text",
                "roleHint": "title",
                "pos": None,
                "paragraphs": [{"text": lines[0], "level": 0,
                                "bullet": False, "size": None, "bold": True}],
                "text": lines[0],
            })
        if len(lines) > 1:
            shapes.append({
                "type": "text",
                "roleHint": "body",
                "pos": None,
                "paragraphs": [{"text": ln, "level": 0, "bullet": False,
                                "size": None, "bold": False}
                               for ln in lines[1:]],
                "text": "\n".join(lines[1:]),
            })
        digest["slides"].append({
            "index": i, "layout": None, "shapes": shapes})
    return digest


# ---------------------------------------------------------------- Markdown

def to_markdown(digest):
    """人が読んで確かめるための digest.md。JSON と同じ内容の読みやすい版。"""
    out = ["# 抽出結果: %s" % digest["source"],
           "",
           "- 形式: %s" % digest["kind"],
           "- スライド数: %d" % len(digest["slides"]),
           ""]
    for slide in digest["slides"]:
        head = "## スライド %d" % slide["index"]
        if slide.get("layout"):
            head += "（レイアウト: %s）" % slide["layout"]
        out += [head, ""]
        for entry in slide["shapes"]:
            role = entry.get("roleHint", "?")
            kind = entry["type"]
            if kind == "table":
                out.append("- **[表]**")
                for row in entry["rows"]:
                    out.append("  | " + " | ".join(row) + " |")
            elif kind == "picture":
                out.append("- **[画像/%s]** `%s` (%sx%s)" % (
                    entry.get("hint", "?"), entry.get("file"),
                    entry.get("w"), entry.get("h")))
            elif kind == "chart":
                out.append("- **[グラフ]** 系列 %d" % len(entry.get("series", [])))
            elif kind == "smartart":
                nodes = entry.get("nodes") or []
                out.append("- **[SmartArt]** ノード %d" % len(nodes))
                for node in nodes:
                    for k, line in enumerate(node.split("\n")):
                        out.append("  %s%s" % ("・" if k == 0 else "  ", line))
                out.append("  > %s" % entry.get("note", ""))
            elif kind == "unsupported":
                out.append("- **[未対応: %s]**" % entry.get("graphic", "?"))
                out.append("  > %s" % entry.get("note", ""))
            else:
                label = "[%s]" % role
                for j, para in enumerate(entry.get("paragraphs") or []):
                    mark = "・" if para["bullet"] else ""
                    indent = "  " * para["level"]
                    prefix = "- **%s** " % label if j == 0 else "  "
                    out.append("%s%s%s%s" % (prefix, indent, mark, para["text"]))
                if entry.get("shapeKind"):
                    out.append("  （図形: %s / 塗り: %s）" % (
                        entry["shapeKind"], entry.get("fill")))
        if slide.get("notes"):
            out += ["", "> ノート: %s" % slide["notes"]]
        out.append("")
    return "\n".join(out)


# ---------------------------------------------------------------- CLI

def main():
    ap = argparse.ArgumentParser(
        description="既存スライド(.pptx/.pdf)から digest.json を抽出する")
    ap.add_argument("input", help="入力ファイル (.pptx または .pdf)")
    ap.add_argument("-o", "--output", default="digest.json",
                    help="出力 JSON (既定: digest.json)")
    ap.add_argument("--media-dir", default=None,
                    help="画像の書き出し先 (既定: 出力JSONと同じ場所の media/)")
    ap.add_argument("--md", default=None,
                    help="人が読む用の Markdown も書き出す")
    args = ap.parse_args()

    src = args.input
    if not os.path.exists(src):
        raise SystemExit("エラー: ファイルが見つかりません → %s" % src)

    ext = os.path.splitext(src)[1].lower()
    out_dir = os.path.dirname(os.path.abspath(args.output))

    if ext == ".pptx":
        media_dir = args.media_dir or os.path.join(out_dir, "media")
        digest = extract_pptx(src, media_dir)
    elif ext == ".pdf":
        digest = extract_pdf(src)
    else:
        raise SystemExit(
            "エラー: 対応していない形式です → %s\n"
            "  .pptx または .pdf を指定してください。\n"
            "  .ppt は PowerPoint で .pptx に保存し直してください。" % ext)

    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
    with open(args.output, "w", encoding="utf-8") as fh:
        json.dump(digest, fh, ensure_ascii=False, indent=2)

    if args.md:
        with open(args.md, "w", encoding="utf-8") as fh:
            fh.write(to_markdown(digest))

    def _count(kind):
        return sum(1 for s in digest["slides"] for e in s["shapes"]
                   if e["type"] == kind)

    msg = "OK: %s  slides=%d  画像=%d" % (
        args.output, len(digest["slides"]), _count("picture"))
    smart, unsup = _count("smartart"), _count("unsupported")
    if smart:
        msg += "  SmartArt=%d" % smart
    if unsup:
        msg += "  未対応=%d" % unsup
    print(msg)
    if smart or unsup:
        print("  ※ SmartArt/未対応の図があります。digest.md の該当箇所を確認し、"
              "figure で描き直してください。")


if __name__ == "__main__":
    main()
